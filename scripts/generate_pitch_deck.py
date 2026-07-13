#!/usr/bin/env python3
"""Generate ClaimGuard-AI investor pitch deck — polished, no competition slides."""

from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    import subprocess
    import sys
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

OUT = Path(__file__).resolve().parent.parent / "ClaimGuard-AI-Pitch-Deck.pptx"

# Brand palette
NAVY   = RGBColor(10,  58,  110)
TEAL   = RGBColor(13, 148, 136)
WHITE  = RGBColor(255, 255, 255)
DARK   = RGBColor(15,  23,  42)
MUTED  = RGBColor(100, 116, 139)
RED    = RGBColor(220, 38,  38)
LIGHT  = RGBColor(241, 245, 249)
ACCENT = RGBColor(224, 242, 254)

W = Inches(13.333)
H = Inches(7.5)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def textbox(slide, l, t, w, h, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tb


def bullets(slide, l, t, w, h, items, size=20, color=DARK, spacing=Pt(10)):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = spacing
        p.level = 0
    return tb


def stat_card(slide, l, t, w, h, number, label, color=NAVY):
    rect(slide, l, t, w, h, LIGHT)
    textbox(slide, l + Inches(0.2), t + Inches(0.15), w - Inches(0.4), Inches(0.7),
            number, size=36, color=color, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, l + Inches(0.1), t + Inches(0.85), w - Inches(0.2), Inches(0.5),
            label, size=13, color=MUTED, align=PP_ALIGN.CENTER)


def slide_header(slide, title, subtitle=None):
    rect(slide, 0, 0, W, Inches(1.35), NAVY)
    textbox(slide, Inches(0.6), Inches(0.25), Inches(12), Inches(0.7),
            title, size=34, color=WHITE, bold=True)
    if subtitle:
        textbox(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
                subtitle, size=16, color=RGBColor(186, 230, 253))


def build(prs):
    # ── SLIDE 1: Title ──────────────────────────────────────────────
    s = blank(prs)
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, Inches(5.8), W, Inches(1.7), TEAL)
    textbox(s, Inches(0.8), Inches(1.8), Inches(11), Inches(1.2),
            "ClaimGuard-AI", size=54, color=WHITE, bold=True)
    textbox(s, Inches(0.8), Inches(3.0), Inches(11), Inches(0.8),
            "Stop insurance denials before the claim leaves the building.", size=24, color=RGBColor(186, 230, 253))
    textbox(s, Inches(0.8), Inches(6.0), Inches(11), Inches(0.5),
            "Priyansh Shah  ·  AIxBio Hackathon @ Bayer Co.Lab  ·  May 2026", size=14, color=WHITE)

    # ── SLIDE 2: Problem ──────────────────────────────────────────
    s = blank(prs)
    slide_header(s, "The Problem", "Healthcare revenue is leaking at the source")
    stat_card(s, Inches(0.6),  Inches(1.7), Inches(3.8), Inches(1.4), "5–10%", "Net revenue lost to denials", RED)
    stat_card(s, Inches(4.7),  Inches(1.7), Inches(3.8), Inches(1.4), "$25", "Cost to rework each denied claim", RED)
    stat_card(s, Inches(8.8),  Inches(1.7), Inches(3.8), Inches(1.4), "42 days", "Average delay from slow payers", RED)
    bullets(s, Inches(0.6), Inches(3.4), Inches(12), Inches(3.5), [
        "Billers review claims in the order they arrive — not by financial impact",
        "Rule-based scrubbers check codes and fields, but never read the physician's note",
        "A missing sentence in a clinical note can trigger a $187,000 oncology denial",
        "By the time a denial letter arrives, the revenue is already gone",
    ], size=20)

    # ── SLIDE 3: Solution ─────────────────────────────────────────
    s = blank(prs)
    slide_header(s, "The Solution", "Three layers of intelligence, one pre-submission workflow")
    layers = [
        ("01  READ",  "Agentic AI parses physician notes", "Flags documentation gaps, CPT/ICD mismatches, missing time for E/M codes"),
        ("02  SCORE", "ML predicts denial probability Pᵢ", "XGBoost model trained on payer + clinical feature signals"),
        ("03  PRIORITIZE", "Financial optimizer ranks the queue", "Expected Loss ELᵢ = Claim Value × Pᵢ  →  auditors fix highest-impact claims first"),
    ]
    for i, (head, sub, detail) in enumerate(layers):
        y = Inches(1.7) + Inches(i * 1.75)
        rect(s, Inches(0.6), y, Inches(12), Inches(1.55), LIGHT if i % 2 == 0 else WHITE)
        rect(s, Inches(0.6), y, Inches(0.12), Inches(1.55), TEAL)
        textbox(s, Inches(0.9), y + Inches(0.12), Inches(3), Inches(0.45), head, size=18, color=TEAL, bold=True)
        textbox(s, Inches(0.9), y + Inches(0.55), Inches(5), Inches(0.4), sub, size=20, color=DARK, bold=True)
        textbox(s, Inches(6.2), y + Inches(0.35), Inches(6.2), Inches(0.8), detail, size=16, color=MUTED)

    # ── SLIDE 4: How it works (architecture) ────────────────────────
    s = blank(prs)
    slide_header(s, "How It Works", "From physician note to prioritized worklist in seconds")
    steps = [
        ("Physician Note\n+ CPT / ICD-10", NAVY),
        ("Groq AI Agent\nStrict JSON extraction", TEAL),
        ("XGBoost\nDenial probability Pᵢ", NAVY),
        ("DuckDB + Knapsack\nSort by Expected Loss", TEAL),
        ("Auditor\nFixes & submits clean claim", NAVY),
    ]
    for i, (label, color) in enumerate(steps):
        x = Inches(0.5 + i * 2.55)
        rect(s, x, Inches(2.5), Inches(2.2), Inches(1.6), color)
        textbox(s, x + Inches(0.1), Inches(2.7), Inches(2.0), Inches(1.2),
                label, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            textbox(s, x + Inches(2.2), Inches(3.1), Inches(0.35), Inches(0.4),
                    "→", size=22, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    bullets(s, Inches(0.6), Inches(4.5), Inches(12), Inches(2.5), [
        "Agent enforces CO-11, CO-16, CO-50, CO-97 denial rules against the clinical narrative",
        "Auditors see the original note, AI reasoning, and a ready-to-paste correction side by side",
        "Treasury Mode re-ranks by cash-flow urgency — slow payers (BCBS 42d) jump the queue",
    ], size=18)

    # ── SLIDE 5: Live Demo cue ──────────────────────────────────────
    s = blank(prs)
    rect(s, 0, 0, W, H, RGBColor(8, 40, 80))
    textbox(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
            "LIVE DEMO", size=20, color=TEAL, bold=True)
    textbox(s, Inches(0.8), Inches(2.1), Inches(11), Inches(1.0),
            "Watch a denial get prevented\nin real time.", size=44, color=WHITE, bold=True)
    demo_steps = [
        "1.  Dashboard  →  Load Pitch Demo",
        "2.  Agent Studio  →  Run analysis on the 99214 note (missing time)",
        "3.  Add one sentence  →  Re-run  →  denial flags disappear",
        "4.  Claims Queue  →  Show how high-value claims float to the top",
    ]
    bullets(s, Inches(0.8), Inches(3.5), Inches(11), Inches(3), demo_steps, size=22, color=RGBColor(186, 230, 253), spacing=Pt(16))
    textbox(s, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
            "→  Switch to browser: http://localhost:3000/dashboard", size=16, color=MUTED)

    # ── SLIDE 6: Product screens ──────────────────────────────────
    s = blank(prs)
    slide_header(s, "The Product", "Four views built for every stakeholder")
    views = [
        ("Dashboard", "Pipeline liquidity, revenue at risk, one-click demo load"),
        ("Claims Queue", "Knapsack-prioritized worklist with agentic flags per claim"),
        ("Agent Studio", "Live note analysis — watch risk score update as you type"),
        ("Reports", "Denial code breakdown, payer trends, PDF/CSV export"),
    ]
    for i, (title, desc) in enumerate(views):
        col, row = i % 2, i // 2
        x = Inches(0.6 + col * 6.3)
        y = Inches(1.7 + row * 2.6)
        rect(s, x, y, Inches(6.0), Inches(2.3), LIGHT)
        rect(s, x, y, Inches(6.0), Inches(0.55), NAVY)
        textbox(s, x + Inches(0.2), y + Inches(0.08), Inches(5.6), Inches(0.4),
                title, size=18, color=WHITE, bold=True)
        textbox(s, x + Inches(0.2), y + Inches(0.75), Inches(5.6), Inches(1.3),
                desc, size=17, color=DARK)

    # ── SLIDE 7: Impact ───────────────────────────────────────────
    s = blank(prs)
    slide_header(s, "Impact", "What a 200-bed health system recovers")
    stat_card(s, Inches(0.6),  Inches(1.7), Inches(3.8), Inches(1.4), "$2.4M", "Annual denial prevention (est.)", TEAL)
    stat_card(s, Inches(4.7),  Inches(1.7), Inches(3.8), Inches(1.4), "+19%", "Cash collections (Treasury mode)", TEAL)
    stat_card(s, Inches(8.8),  Inches(1.7), Inches(3.8), Inches(1.4), "8/day", "Claims an auditor can actually fix", TEAL)
    bullets(s, Inches(0.6), Inches(3.4), Inches(12), Inches(3.5), [
        "Pre-submission prevention eliminates the $25 rework cost per denial entirely",
        "Bounded knapsack ensures staff capacity goes to highest Expected Loss claims",
        "Agent generates ready-to-paste clinical corrections — not just error codes",
        "Policy check and auto-appeals built in for the claims that still slip through",
    ], size=20)

    # ── SLIDE 8: Business model ───────────────────────────────────
    s = blank(prs)
    slide_header(s, "Business Model", "SaaS + performance — aligned with customer ROI")
    rect(s, Inches(0.6), Inches(1.7), Inches(5.8), Inches(4.8), LIGHT)
    textbox(s, Inches(0.9), Inches(1.9), Inches(5.2), Inches(0.5),
            "Tier 1 — Platform SaaS", size=22, color=NAVY, bold=True)
    bullets(s, Inches(0.9), Inches(2.5), Inches(5.2), Inches(3.5), [
        "$0.15 – $0.40 per claim monitored",
        "Scales with claim volume",
        "Includes agent, ML scoring, queue",
    ], size=18)
    rect(s, Inches(6.8), Inches(1.7), Inches(5.8), Inches(4.8), ACCENT)
    textbox(s, Inches(7.1), Inches(1.9), Inches(5.2), Inches(0.5),
            "Tier 2 — Performance Fee", size=22, color=NAVY, bold=True)
    bullets(s, Inches(7.1), Inches(2.5), Inches(5.2), Inches(3.5), [
        "5–10% of net cash flow unlocked",
        "Measured in first 90-day pilot",
        "Aligns our success with theirs",
    ], size=18)
    textbox(s, Inches(0.6), Inches(6.65), Inches(12), Inches(0.5),
            "Target customer: 200-bed regional health system  ·  $50M–$500M net patient revenue",
            size=15, color=MUTED)

    # ── SLIDE 9: Traction ─────────────────────────────────────────
    s = blank(prs)
    slide_header(s, "Where We Are Today", "Working product, ready for pilot")
    bullets(s, Inches(0.6), Inches(1.7), Inches(6), Inches(5), [
        "✓  Full-stack MVP deployed locally",
        "✓  Groq agent with strict JSON schema",
        "✓  XGBoost denial probability engine",
        "✓  DuckDB analytics + knapsack optimizer",
        "✓  6 realistic demo claims (UHC, Aetna, Medicare, BCBS, Medicaid)",
        "✓  Policy check + auto-appeals live",
    ], size=20)
    rect(s, Inches(7.2), Inches(1.7), Inches(5.5), Inches(5), NAVY)
    textbox(s, Inches(7.5), Inches(2.0), Inches(5), Inches(0.5),
            "Next 90 Days", size=22, color=TEAL, bold=True)
    bullets(s, Inches(7.5), Inches(2.7), Inches(5), Inches(3.5), [
        "Design partner health system",
        "Pilot on top 5 denial CPT codes",
        "Measure $ recovered vs. control group",
        "Deploy to Vercel + Render",
    ], size=18, color=WHITE)

    # ── SLIDE 10: Ask / Close ───────────────────────────────────────
    s = blank(prs)
    rect(s, 0, 0, W, H, NAVY)
    textbox(s, Inches(0.8), Inches(2.0), Inches(11), Inches(0.6),
            "THE ASK", size=20, color=TEAL, bold=True)
    textbox(s, Inches(0.8), Inches(2.7), Inches(11), Inches(1.2),
            "One design partner.\nOne 90-day pilot.", size=44, color=WHITE, bold=True)
    bullets(s, Inches(0.8), Inches(4.3), Inches(11), Inches(2.5), [
        "Intro to a revenue cycle director at a regional health system",
        "Seed funding to deploy, hire 1 ML engineer, run the pilot",
        "We will run your highest-denial CPT codes live — today",
    ], size=22, color=RGBColor(186, 230, 253), spacing=Pt(14))
    textbox(s, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
            "priyansh.shah@stonybrook.edu  ·  ClaimGuard-AI", size=16, color=MUTED)


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    build(prs)
    prs.save(OUT)
    print(f"Created: {OUT}")


if __name__ == "__main__":
    main()
