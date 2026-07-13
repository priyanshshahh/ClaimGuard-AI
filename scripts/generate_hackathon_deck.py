#!/usr/bin/env python3
"""Generate 5-slide hackathon lightning pitch deck (AIxBio @ Bayer Co.Lab)."""

from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    import subprocess
    import sys
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

OUT = Path(__file__).resolve().parent.parent / "ClaimGuard-AI-Lightning-Pitch.pptx"

NAVY = RGBColor(10, 58, 110)
TEAL = RGBColor(13, 148, 136)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(15, 23, 42)
MUTED = RGBColor(100, 116, 139)
RED = RGBColor(220, 38, 38)
LIGHT = RGBColor(241, 245, 249)
W = Inches(13.333)
H = Inches(7.5)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()


def tb(slide, l, t, w, h, text, size=20, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return box


def bullets(slide, l, t, w, h, items, size=22, color=DARK):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(12)


def header(slide, title, subtitle=None):
    rect(slide, 0, 0, W, Inches(1.3), NAVY)
    tb(slide, Inches(0.6), Inches(0.22), Inches(12), Inches(0.65), title, 34, WHITE, True)
    if subtitle:
        tb(slide, Inches(0.6), Inches(0.82), Inches(12), Inches(0.4), subtitle, 16, RGBColor(186, 230, 253))


def stat(slide, l, t, num, label, color=RED):
    rect(slide, l, t, Inches(3.8), Inches(1.35), LIGHT)
    tb(slide, l + Inches(0.15), t + Inches(0.15), Inches(3.5), Inches(0.65), num, 38, color, True, PP_ALIGN.CENTER)
    tb(slide, l + Inches(0.1), t + Inches(0.82), Inches(3.6), Inches(0.45), label, 13, MUTED, False, PP_ALIGN.CENTER)


def build(prs):
    # SLIDE 1 — Title + Problem hook
    s = blank(prs)
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, Inches(6.2), W, Inches(1.3), TEAL)
    tb(s, Inches(0.7), Inches(1.5), Inches(11), Inches(0.5), "AIxBio Hackathon · Bayer Co.Lab · May 2026", 14, RGBColor(186, 230, 253))
    tb(s, Inches(0.7), Inches(2.1), Inches(11), Inches(1.1), "ClaimGuard-AI", 52, WHITE, True)
    tb(s, Inches(0.7), Inches(3.3), Inches(11), Inches(0.9),
       "The Problem: Hospitals lose 5–10% of net revenue to preventable claim denials.", 26, WHITE)
    bullets(s, Inches(0.7), Inches(4.3), Inches(11), Inches(1.8), [
        "Each denied claim costs ~$25 to rework",
        "Manual auditing cannot scale",
        "Scrubbers never read the physician's note",
    ], 20, RGBColor(186, 230, 253))
    tb(s, Inches(0.7), Inches(6.45), Inches(11), Inches(0.4), "Priyansh Shah · Stony Brook University", 14, WHITE)

    # SLIDE 2 — Solution
    s = blank(prs)
    header(s, "The Solution", "ClaimGuard-AI — pre-submission denial prevention")
    bullets(s, Inches(0.7), Inches(1.6), Inches(11.5), Inches(5), [
        "Agentic AI reads unstructured physician notes via Nebius Token Factory (Gemma 3 27B)",
        "Strict JSON extraction: documentation gaps, CPT/ICD mismatches, denial codes",
        "XGBoost predicts denial probability Pᵢ for every claim",
        "Bounded knapsack sorts queue by Expected Financial Loss: ELᵢ = Vᵢ × Pᵢ",
        "Billers fix high-impact claims first — before submission, not after denial",
    ], 24)
    rect(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.9), LIGHT)
    tb(s, Inches(0.9), Inches(6.15), Inches(11.5), Inches(0.6),
       "Full-stack: Next.js frontend · FastAPI backend · DuckDB analytics · Render + Vercel ready", 16, TEAL, True)

    # SLIDE 3 — 24-Hour Progress
    s = blank(prs)
    header(s, "24-Hour Progress", "What we built at the hackathon")
    stat(s, Inches(0.6), Inches(1.65), "6", "Demo claims analyzed live")
    stat(s, Inches(4.7), Inches(1.65), "3", "Product pages shipped", TEAL)
    stat(s, Inches(8.8), Inches(1.65), "4", "AI pipeline layers", TEAL)
    bullets(s, Inches(0.7), Inches(3.3), Inches(5.5), Inches(3.5), [
        "✓  Multi-page React + FastAPI app",
        "✓  Nebius API integration (strict JSON)",
        "✓  XGBoost denial classifier",
        "✓  DuckDB + knapsack optimizer",
        "✓  Policy check + auto-appeals",
        "✓  Executive dashboard + exports",
    ], 19)
    rect(s, Inches(6.8), Inches(3.3), Inches(5.8), Inches(3.5), NAVY)
    tb(s, Inches(7.1), Inches(3.5), Inches(5.2), Inches(0.45), "Live Demo (90 sec)", 20, TEAL, True)
    bullets(s, Inches(7.1), Inches(4.1), Inches(5.2), Inches(2.5), [
        "Load 6 hospital claims",
        "99214 note → add time → denial flags clear",
        "Treasury queue re-ranks by cash urgency",
    ], 17, WHITE)

    # SLIDE 4 — Future Directions
    s = blank(prs)
    header(s, "Future Directions", "From auditing co-pilot to autonomous clearinghouse")
    bullets(s, Inches(0.7), Inches(1.6), Inches(11.5), Inches(4.5), [
        "Pilot with a regional health system — measure $ recovered in 90 days",
        "Expand medical necessity rules across all specialties",
        "Supabase persistence + EHR integration (FHIR R4)",
        "Evolve knapsack optimizer into fully autonomous financial clearinghouse",
        "Global Sprint Hackathon (November) — production deployment at scale",
    ], 24)
    tb(s, Inches(0.7), Inches(6.2), Inches(11), Inches(0.6),
       "Ask: Design partner health system + Global Sprint selection", 20, TEAL, True)

    # SLIDE 5 — Close / Live demo cue
    s = blank(prs)
    rect(s, 0, 0, W, H, TEAL)
    tb(s, Inches(0.8), Inches(2.2), Inches(11), Inches(1.2),
       "Let's prevent a denial\nlive.", 48, WHITE, True, PP_ALIGN.CENTER)
    tb(s, Inches(0.8), Inches(4.0), Inches(11), Inches(0.6),
       "localhost:3000/dashboard → Load Pitch Demo → Agent Studio", 22, NAVY, False, PP_ALIGN.CENTER)
    tb(s, Inches(0.8), Inches(5.5), Inches(11), Inches(0.5),
       "priyansh.shah@stonybrook.edu", 18, WHITE, False, PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    build(prs)
    prs.save(OUT)
    print(f"Created: {OUT}")


if __name__ == "__main__":
    main()
